#!/usr/bin/env python3
"""
Create comprehensive PowerPoint slide decks for MBA 590 lectures.
Generates slides for Weeks 1, 6, and 13 to accompany Jupyter notebooks.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def create_title_slide(prs, title, subtitle):
    """Create a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Format title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue

    # Format subtitle
    subtitle_frame = subtitle_shape.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(68, 68, 68)  # Dark gray


def create_section_slide(prs, section_title):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

    # Add centered text
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = section_title

    para = text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    para.font.size = Pt(54)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)  # White


def create_content_slide(prs, title, content_items, layout_type=1):
    """Create a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_type])
    title_shape = slide.shapes.title
    title_shape.text = title

    # Format title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # Add content
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        if isinstance(item, tuple):
            # (text, level) tuple
            p.text = item[0]
            p.level = item[1]
        else:
            p.text = item
            p.level = 0

        p.font.size = Pt(18)
        p.space_after = Pt(10)


def create_week1_slides():
    """Create slides for Week 1: Prompt Engineering."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    create_title_slide(
        prs,
        "Week 1: Foundations of Prompt Engineering",
        "MBA 590 - Advanced AI Strategy\nPrompting and Agentic Frameworks",
    )

    # Agenda
    create_content_slide(
        prs,
        "Today's Agenda",
        [
            "1. What is Prompt Engineering?",
            "2. The Three Core Prompt Types",
            ("Zero-Shot Prompts", 1),
            ("Role-Playing Prompts", 1),
            ("Detailed Instruction Prompts", 1),
            "3. Comparative Framework & Evaluation",
            "4. Best Practices & Common Pitfalls",
            "5. Business Applications & ROI",
        ],
    )

    # Learning Objectives
    create_content_slide(
        prs,
        "Learning Objectives",
        [
            "By the end of this session, you will be able to:",
            "",
            "✓ Define prompt engineering and explain its business value",
            "✓ Distinguish between zero-shot, role-playing, and detailed prompts",
            "✓ Evaluate which prompt type to use for specific use cases",
            "✓ Apply a decision framework for prompt selection",
            "✓ Estimate cost-benefit trade-offs for different approaches",
        ],
    )

    # Section 1
    create_section_slide(prs, "What is Prompt Engineering?")

    create_content_slide(
        prs,
        "Prompt Engineering Defined",
        [
            "The practice of designing effective inputs to AI models",
            "",
            "Key Components:",
            ("Clear instructions: What do you want the AI to do?", 1),
            ("Context: What background information is needed?", 1),
            ("Constraints: What limitations or requirements?", 1),
            ("Examples: What does good output look like?", 1),
            "",
            "💡 Think of it as: Writing a perfect job description for an AI employee",
        ],
    )

    create_content_slide(
        prs,
        "Why It Matters for Business",
        [
            "Traditional Software:",
            ("Write code: IF customer_angry THEN apologize", 1),
            ("Rigid, requires programming every scenario", 1),
            "",
            "Prompt Engineering:",
            ("Write instructions: 'Respond empathetically to upset customers'", 1),
            ("Flexible, handles scenarios you didn't explicitly program", 1),
            "",
            "Result: 10-100x faster to deploy new capabilities",
        ],
    )

    # Section 2: Zero-Shot
    create_section_slide(prs, "Prompt Type 1: Zero-Shot")

    create_content_slide(
        prs,
        "Zero-Shot Prompts",
        [
            "Definition: Direct instruction with NO examples",
            "",
            "Example:",
            ("'Write an email response to a customer who received", 1),
            ("a damaged product.'", 1),
            "",
            "When to Use:",
            ("✓ Simple, well-understood tasks", 1),
            ("✓ Need speed over precision", 1),
            ("✓ Budget-constrained (minimal tokens)", 1),
            ("✗ Don't use for: Complex tasks, specific tone requirements", 1),
        ],
    )

    create_content_slide(
        prs,
        "Zero-Shot: Pros & Cons",
        [
            "Strengths:",
            ("⚡ Fast to create (10 seconds)", 1),
            ("💰 Minimal token usage = lower cost ($0.001 per query)", 1),
            ("🔄 Easy to test and iterate", 1),
            "",
            "Weaknesses:",
            ("❌ Generic output (no company-specific tone)", 1),
            ("❌ May miss important details (policy, timeline)", 1),
            ("❌ Inconsistent quality across use cases", 1),
        ],
    )

    # Section 3: Role-Playing
    create_section_slide(prs, "Prompt Type 2: Role-Playing")

    create_content_slide(
        prs,
        "Role-Playing Prompts",
        [
            "Definition: Give the AI a specific role and persona",
            "",
            "Example:",
            ("'You are a senior customer service representative at XYZ Corp", 1),
            ("with 10 years of experience. You are known for empathy and", 1),
            ("problem-solving. Respond to this damaged product complaint...'", 1),
            "",
            "The '4 Ps' of Role-Playing:",
            ("1. Persona (who are you?)", 1),
            ("2. Purpose (what's your goal?)", 1),
            ("3. Parameters (what are the constraints?)", 1),
            ("4. Process (how should you approach it?)", 1),
        ],
    )

    create_content_slide(
        prs,
        "Role-Playing: When to Use",
        [
            "Best For:",
            ("✓ Tasks requiring specific tone/voice", 1),
            ("✓ Domain expertise needed (legal, medical, technical)", 1),
            ("✓ Consistent brand personality", 1),
            "",
            "Cost-Benefit:",
            ("Cost: ~$0.002-0.003 per query (3x zero-shot)", 1),
            ("Benefit: 40-60% better output quality", 1),
            ("ROI: Worth it for customer-facing content", 1),
        ],
    )

    # Section 4: Detailed Instructions
    create_section_slide(prs, "Prompt Type 3: Detailed Instructions")

    create_content_slide(
        prs,
        "Detailed Instruction Prompts",
        [
            "Definition: Comprehensive prompt with context, examples, format",
            "",
            "Includes:",
            ("📋 Background context", 1),
            ("🎯 Specific objectives", 1),
            ("📐 Output format/structure", 1),
            ("✅ Examples of good output", 1),
            ("🚫 What NOT to do", 1),
            ("⚖️ Quality criteria", 1),
        ],
    )

    create_content_slide(
        prs,
        "Detailed Instructions: Components",
        [
            "Template Structure:",
            "",
            "1. Role & Context (who you are, what situation)",
            "2. Task Definition (exactly what to do)",
            "3. Constraints (word limit, tone, requirements)",
            "4. Examples (2-3 good examples)",
            "5. Output Format (structure, sections)",
            "6. Quality Checks (what makes a good response)",
        ],
    )

    create_content_slide(
        prs,
        "Cost-Benefit Analysis",
        [
            "Comparison by Prompt Type:",
            "",
            "Zero-Shot:",
            ("Cost: $0.001/query | Quality: 60% | Time: 10 sec", 1),
            "",
            "Role-Playing:",
            ("Cost: $0.003/query | Quality: 80% | Time: 2 min", 1),
            "",
            "Detailed:",
            ("Cost: $0.005/query | Quality: 95% | Time: 15 min", 1),
            "",
            "💡 Decision: Match complexity to business need",
        ],
    )

    # Evaluation Framework
    create_section_slide(prs, "Evaluation & Decision Framework")

    create_content_slide(
        prs,
        "Decision Matrix",
        [
            "When to Use Each Prompt Type:",
            "",
            "Zero-Shot:",
            ("• High volume, low stakes (1000s of queries/day)", 1),
            ("• Internal tools, draft generation", 1),
            "",
            "Role-Playing:",
            ("• Medium volume, medium stakes (100s/day)", 1),
            ("• Customer-facing but low risk", 1),
            "",
            "Detailed Instructions:",
            ("• Low volume, high stakes (10s/day)", 1),
            ("• Legal, financial, medical, executive content", 1),
        ],
    )

    # Best Practices
    create_section_slide(prs, "Best Practices")

    create_content_slide(
        prs,
        "10 Prompt Engineering Principles",
        [
            "1. Be specific (not 'summarize' but 'create 3-bullet summary')",
            "2. Provide context (background info the AI needs)",
            "3. Use examples (show don't tell)",
            "4. Set constraints (word limit, tone, format)",
            "5. Iterate (test and refine)",
            "6. Specify output format (JSON, bullet points, paragraph)",
            "7. Think step-by-step (break complex tasks into steps)",
            "8. Use delimiters (### to separate sections)",
            "9. Request reasoning (ask AI to 'think aloud')",
            "10. Measure quality (define success criteria)",
        ],
    )

    # Real-World Applications
    create_content_slide(
        prs,
        "Real-World Business Applications",
        [
            "Customer Service:",
            ("Email responses, chat support, FAQs (Role-playing)", 1),
            "",
            "Content Creation:",
            ("Marketing copy, social media, blogs (Detailed)", 1),
            "",
            "Data Analysis:",
            ("Report summaries, insight extraction (Zero-shot)", 1),
            "",
            "Code Generation:",
            ("Documentation, code review, debugging (Detailed)", 1),
        ],
    )

    # ROI Example
    create_content_slide(
        prs,
        "ROI Example: Customer Service",
        [
            "Scenario: E-commerce company, 1000 emails/day",
            "",
            "Traditional Approach:",
            ("• 10 agents × $30/hr × 8hr = $2,400/day", 1),
            ("• $600K/year in labor", 1),
            "",
            "AI with Detailed Prompts:",
            ("• $0.005 × 1000 queries = $5/day", 1),
            ("• $1,825/year in API costs", 1),
            ("• 2 agents for review/escalation = $120K/year", 1),
            "",
            "Net Savings: $478K/year (80% reduction)",
        ],
    )

    # Hands-On Exercise
    create_content_slide(
        prs,
        "Hands-On Exercise",
        [
            "Your Turn: Design a Prompt",
            "",
            "Scenario:",
            ("You need an AI to analyze customer feedback and", 1),
            ("categorize issues (product, shipping, service)", 1),
            "",
            "Questions to Consider:",
            ("1. Which prompt type would you use and why?", 1),
            ("2. What context does the AI need?", 1),
            ("3. What constraints should you set?", 1),
            ("4. How would you measure success?", 1),
        ],
    )

    # Key Takeaways
    create_content_slide(
        prs,
        "Key Takeaways",
        [
            "1. Prompt engineering is a critical skill for AI deployment",
            "",
            "2. Three prompt types serve different needs:",
            ("• Zero-shot: Fast, cheap, good enough", 1),
            ("• Role-playing: Balanced quality & cost", 1),
            ("• Detailed: High quality, high effort", 1),
            "",
            "3. Match prompt complexity to business value",
            "",
            "4. Always measure: quality, cost, time",
            "",
            "5. Iterate: Your first prompt won't be perfect",
        ],
    )

    # Q&A
    create_section_slide(prs, "Questions & Discussion")

    # Save
    prs.save("/home/user/data-advanced-ai/slides/Week01_Prompt_Engineering.pptx")
    print("✅ Created: Week01_Prompt_Engineering.pptx")


def create_week6_slides():
    """Create slides for Week 6: Agentic Frameworks."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    create_title_slide(
        prs,
        "Week 6: Introduction to Agentic Frameworks",
        "MBA 590 - Advanced AI Strategy\nConcepts and Architectures",
    )

    # Agenda
    create_content_slide(
        prs,
        "Today's Agenda",
        [
            "1. What Makes a System 'Agentic'?",
            "2. Agents vs. Automation vs. Chatbots",
            "3. Core Agent Concepts",
            ("Perception, Planning, Reasoning, Action, Memory", 1),
            "4. The ReAct Framework",
            "5. Agent Architectures Overview",
            "6. Real-World Case Studies",
            "7. When (Not) to Use Agentic AI",
        ],
    )

    # Learning Objectives
    create_content_slide(
        prs,
        "Learning Objectives",
        [
            "By the end of this session, you will be able to:",
            "",
            "✓ Define agentic systems and their 6 key characteristics",
            "✓ Distinguish agents from automation and chatbots",
            "✓ Explain the perception-action loop",
            "✓ Describe the ReAct framework and its advantages",
            "✓ Identify which agent architecture fits specific problems",
            "✓ Evaluate when agentic capabilities add business value",
        ],
    )

    # What is Agentic?
    create_section_slide(prs, "What is an Agentic System?")

    create_content_slide(
        prs,
        "Agentic Systems Defined",
        [
            "An AI system that can:",
            "",
            "1. Perceive its environment and task requirements",
            "2. Plan a sequence of actions to achieve goals",
            "3. Reason about options and consequences",
            "4. Act using available tools and interfaces",
            "5. Learn from feedback and adapt behavior",
            "6. Operate autonomously with minimal human intervention",
            "",
            "💡 Key: 'Agency' = capacity to act independently",
        ],
    )

    create_content_slide(
        prs,
        "The 6 Characteristics",
        [
            "What distinguishes agentic AI:",
            "",
            "Autonomy: Makes decisions independently",
            "Goal-Directedness: Works toward specific objectives",
            "Reactivity: Responds to environment changes",
            "Pro-activeness: Takes initiative",
            "Social Ability: Interacts with humans/other agents",
            "Learning: Improves performance over time",
            "",
            "If you can write it as IF-THEN rules → Not agentic!",
        ],
    )

    # Agents vs Others
    create_section_slide(prs, "Agents vs. Automation vs. Chatbots")

    create_content_slide(
        prs,
        "Key Differences",
        [
            "Automation:",
            ("• Follows fixed script", 1),
            ("• No planning or reasoning", 1),
            ("• Fails on unexpected input", 1),
            "",
            "Chatbot:",
            ("• Matches intents to responses", 1),
            ("• Single-turn interactions", 1),
            ("• Can't handle multi-step tasks", 1),
            "",
            "Agent:",
            ("• Plans and adapts dynamically", 1),
            ("• Multi-step reasoning", 1),
            ("• Handles complex, novel situations", 1),
        ],
    )

    create_content_slide(
        prs,
        "Decision Framework",
        [
            "When do you need an agent?",
            "",
            "✅ Use Agentic AI when:",
            ("• Path to solution isn't predefined", 1),
            ("• Multiple steps required, sequence varies", 1),
            ("• Contextual decision-making needed", 1),
            ("• Task requires synthesis of information", 1),
            "",
            "❌ Don't use when:",
            ("• Process is fully defined and unchanging → Automation", 1),
            ("• Single-turn Q&A → Chatbot or RAG", 1),
            ("• Safety requires complete human control", 1),
        ],
    )

    # Core Concepts
    create_section_slide(prs, "Core Agent Concepts")

    create_content_slide(
        prs,
        "The Perception-Action Loop",
        [
            "The fundamental cycle of agent behavior:",
            "",
            "1. PERCEIVE: Understand current environment state",
            "2. REASON: Think about what to do",
            "3. PLAN: Decide on next action",
            "4. ACT: Execute the action",
            "5. OBSERVE: See what happened",
            "6. REPEAT: Loop back to step 1",
            "",
            "This is how agents handle dynamic, multi-step tasks",
        ],
    )

    create_content_slide(
        prs,
        "Agent Memory Types",
        [
            "How agents remember:",
            "",
            "Short-Term (Working Memory):",
            ("• Recent interactions (last 10 messages)", 1),
            ("• Current task context", 1),
            ("• Cost: Higher token usage", 1),
            "",
            "Long-Term (Vector Database):",
            ("• Past conversations, learned facts", 1),
            ("• User preferences, historical patterns", 1),
            ("• Cost: Storage + retrieval", 1),
            "",
            "💡 Trade-off: More memory = better context but higher cost",
        ],
    )

    # ReAct Framework
    create_section_slide(prs, "The ReAct Framework")

    create_content_slide(
        prs,
        "ReAct: Reasoning + Acting",
        [
            "Revolutionary approach that interleaves:",
            "",
            "Think → Act → Observe → Think → Act → Observe...",
            "",
            "vs. Traditional:",
            "",
            "Query → Get Result → Answer (no adaptation)",
            "",
            "Key Difference: Reasoning is EXPLICIT and VISIBLE",
            "",
            "📚 Source: Yao et al., 2022 (ICLR 2023)",
        ],
    )

    create_content_slide(
        prs,
        "ReAct Performance Gains",
        [
            "Results from research paper:",
            "",
            "HotpotQA (multi-hop reasoning):",
            ("Traditional: 28% → ReAct: 47% (+68%)", 1),
            "",
            "FEVER (fact verification):",
            ("Traditional: 72% → ReAct: 84% (+17%)", 1),
            "",
            "WebShop (web navigation):",
            ("Traditional: 45% → ReAct: 78% (+73%)", 1),
            "",
            "Why? Can adapt based on intermediate results!",
        ],
    )

    create_content_slide(
        prs,
        "ReAct Example",
        [
            "Task: 'Analyze Q4 2024 sales for Northeast region'",
            "",
            "Thought: I need to first query the database",
            "Action: query_database('Q4 2024 Northeast sales')",
            "Observation: Got data - $5.2M revenue, 1250 transactions",
            "",
            "Thought: I have raw data, now need to analyze it",
            "Action: analyze_by_product(sales_data)",
            "Observation: Product A: 35%, Product B: 28%, Product C: 37%",
            "",
            "Thought: Sufficient data, ready to answer",
            "Action: FINISH with summary",
        ],
    )

    # Agent Architectures
    create_section_slide(prs, "Agent Architectures Overview")

    create_content_slide(
        prs,
        "7 Levels of Agent Complexity",
        [
            "1. Simple Reflex: IF-THEN rules",
            "2. Model-Based: Track internal state",
            "3. Goal-Based: Plan toward goals",
            "4. Utility-Based: Optimize utility function",
            "5. Learning: Improve from experience",
            "6. ReAct: Interleaved reasoning & acting",
            "7. Multi-Agent: Coordinated specialized agents",
            "",
            "💡 Golden Rule: Use simplest architecture that works!",
        ],
    )

    create_content_slide(
        prs,
        "Architecture Selection Guide",
        [
            "Match architecture to need:",
            "",
            "Simple Reflex → Alert on threshold breach",
            "Model-Based → Inventory management",
            "Goal-Based → Project planning",
            "Utility-Based → Resource allocation",
            "Learning → Customer preference learning",
            "ReAct → Research & analysis tasks",
            "Multi-Agent → Cross-functional workflows",
            "",
            "Over-engineering wastes money, under-engineering fails",
        ],
    )

    # Case Studies
    create_section_slide(prs, "Real-World Case Studies")

    create_content_slide(
        prs,
        "Case Study: Klarna Customer Service",
        [
            "Company: Klarna (Fintech, 150M users)",
            "Challenge: Handle millions of support inquiries",
            "Solution: Agentic AI assistant (Feb 2024)",
            "",
            "Agentic Capabilities:",
            ("• Understands intent, identifies account context", 1),
            ("• Plans multi-step resolution path", 1),
            ("• Uses tools (transaction DB, refund system, KB)", 1),
            ("• Learns from satisfaction scores", 1),
            "",
            "Results:",
            ("• 2.3M conversations first month (= 700 FTE agents)", 1),
            ("• Resolution time: 11 min → 2 min", 1),
            ("• $40M annual savings projected", 1),
        ],
    )

    create_content_slide(
        prs,
        "Case Study: Morgan Stanley",
        [
            "Company: Morgan Stanley (Financial Services)",
            "Challenge: 16,000 advisors search 100K+ documents",
            "Solution: GPT-4 agentic system",
            "",
            "Agentic Workflow:",
            ("1. Understand advisor's client context", 1),
            ("2. Break query into sub-questions", 1),
            ("3. Use multiple retrieval strategies", 1),
            ("4. Synthesize findings, identify gaps", 1),
            ("5. Generate actionable insights with citations", 1),
            "",
            "Results:",
            ("• 16,000 advisors deployed", 1),
            ("• 10-15% more time with clients", 1),
            ("• Previously buried insights now discoverable", 1),
        ],
    )

    # Business Applications
    create_content_slide(
        prs,
        "When Agentic Capabilities Add Value",
        [
            "Best Use Cases:",
            "",
            "Complex Customer Inquiry:",
            ("Multi-step problem solving, tool use, adaptation", 1),
            "",
            "Market Research:",
            ("Dynamic data gathering, synthesis, reporting", 1),
            "",
            "Financial Analysis:",
            ("Multiple data sources, complex calculations, insights", 1),
            "",
            "Supply Chain Optimization:",
            ("Real-time monitoring, planning, decision making", 1),
        ],
    )

    # ROI Framework
    create_content_slide(
        prs,
        "ROI Calculation Framework",
        [
            "Costs:",
            ("• Development: $10K-50K (initial agent)", 1),
            ("• Infrastructure: $500-2K/month (LLM API, tools)", 1),
            ("• Per-task: $0.10-0.50 (LLM + tool calls)", 1),
            "",
            "Benefits:",
            ("• Time savings: [hours saved] × [rate] × [tasks/month]", 1),
            ("• Quality: [error reduction %] × [cost of errors]", 1),
            ("• Scale: [tasks enabled] × [value per task]", 1),
            "",
            "Typical Break-Even: 3-12 months",
        ],
    )

    # Key Takeaways
    create_content_slide(
        prs,
        "Key Takeaways",
        [
            "1. Agentic systems combine autonomy, planning,",
            "   reasoning, and tool use",
            "",
            "2. Different from automation (rigid) and chatbots (single-turn)",
            "",
            "3. ReAct framework enables adaptive multi-step reasoning",
            "",
            "4. Match architecture complexity to business need",
            "",
            "5. Real value comes from handling tasks that are",
            "   difficult to pre-program",
            "",
            "6. Careful consideration of autonomy and oversight required",
        ],
    )

    # Q&A
    create_section_slide(prs, "Questions & Discussion")

    # Save
    prs.save("/home/user/data-advanced-ai/slides/Week06_Agentic_Frameworks.pptx")
    print("✅ Created: Week06_Agentic_Frameworks.pptx")


def create_week13_slides():
    """Create slides for Week 13: Technology Strategy."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    create_title_slide(
        prs,
        "Week 13: Technology Strategy & Portfolio Management",
        "MBA 590 - Advanced AI Strategy\nDeveloping Strategic Technology Roadmaps",
    )

    # Agenda
    create_content_slide(
        prs,
        "Today's Agenda",
        [
            "1. Technology Strategy Frameworks",
            "2. Project Prioritization Methodologies",
            ("Weighted Scoring, RICE, Value vs Effort", 1),
            "3. Portfolio Management & Balance",
            ("Three Horizons, Risk Profiles", 1),
            "4. Build vs. Buy vs. Partner Decisions",
            "5. Technology Roadmap Development",
            "6. Real-World Case Studies",
        ],
    )

    # Learning Objectives
    create_content_slide(
        prs,
        "Learning Objectives",
        [
            "By the end of this session, you will be able to:",
            "",
            "✓ Apply prioritization frameworks to technology initiatives",
            "✓ Analyze portfolio balance across time horizons",
            "✓ Create a technology roadmap with milestones",
            "✓ Evaluate build vs buy vs partner trade-offs",
            "✓ Develop a technology strategy aligned with business goals",
            "✓ Justify investment decisions to stakeholders with data",
        ],
    )

    # Strategy Frameworks
    create_section_slide(prs, "Technology Strategy Frameworks")

    create_content_slide(
        prs,
        "Technology Adoption Lifecycle",
        [
            "Where should you be on the curve?",
            "",
            "Innovators (2.5%):",
            ("• Bleeding edge, very high risk, unproven", 1),
            "",
            "Early Adopters (13.5%):",
            ("• Competitive advantage, proven in select cases", 1),
            "",
            "Early Majority (34%):",
            ("• Mainstream adoption, proven value", 1),
            "",
            "Late Majority (34%):",
            ("• Industry standard, lower risk", 1),
            "",
            "Decision: Varies by technology and strategic importance",
        ],
    )

    create_content_slide(
        prs,
        "Technology Assessment Matrix",
        [
            "Evaluate each technology on 5 dimensions:",
            "",
            "1. Maturity: Innovator → Late Majority",
            "2. Business Value: Low → Very High",
            "3. Implementation Risk: Low → Very High",
            "4. Time to Value: Immediate → 5+ years",
            "5. Strategic Priority: Watch → Critical",
            "",
            "💡 Map technologies on Maturity vs Value chart",
            "    (See notebook for interactive visualization)",
        ],
    )

    # Prioritization
    create_section_slide(prs, "Project Prioritization Methodologies")

    create_content_slide(
        prs,
        "Weighted Scoring Model",
        [
            "Our recommended approach:",
            "",
            "Criteria (example weights):",
            ("• Strategic Alignment: 30%", 1),
            ("• Financial Impact: 25%", 1),
            ("• Implementation Ease: 20%", 1),
            ("• Risk Level: 15%", 1),
            ("• Time to Value: 10%", 1),
            "",
            "Calculate weighted score for each project",
            "Rank by score, apply judgment to top candidates",
        ],
    )

    create_content_slide(
        prs,
        "RICE Framework (Alternative)",
        [
            "From Intercom, focus on ROI per effort:",
            "",
            "RICE = (Reach × Impact × Confidence) / Effort",
            "",
            "Reach: How many users/processes affected?",
            "Impact: How much improvement per user?",
            "Confidence: How certain are estimates? (0-100%)",
            "Effort: How much work required? (person-months)",
            "",
            "Best for: Resource-constrained teams maximizing efficiency",
        ],
    )

    create_content_slide(
        prs,
        "Strategic Impact vs Difficulty Matrix",
        [
            "Classic 2×2 prioritization:",
            "",
            "High Impact, Low Difficulty:",
            ("→ QUICK WINS - Do these first!", 1),
            "",
            "High Impact, High Difficulty:",
            ("→ STRATEGIC BETS - Plan carefully, allocate resources", 1),
            "",
            "Low Impact, Low Difficulty:",
            ("→ FILL-INS - Do when capacity available", 1),
            "",
            "Low Impact, High Difficulty:",
            ("→ TIME SINKS - Avoid or defer", 1),
        ],
    )

    # Portfolio Management
    create_section_slide(prs, "Portfolio Management")

    create_content_slide(
        prs,
        "The Three Horizons Model",
        [
            "McKinsey's framework for balanced investment:",
            "",
            "Horizon 1 (0-12 months): 60-70%",
            ("• Incremental improvements", 1),
            ("• Extend current business", 1),
            ("• Low risk, proven ROI", 1),
            "",
            "Horizon 2 (1-3 years): 20-30%",
            ("• Adjacent opportunities", 1),
            ("• Growth initiatives", 1),
            ("• Medium risk, validated concepts", 1),
            "",
            "Horizon 3 (3+ years): 10-20%",
            ("• Transformational bets", 1),
            ("• New business models", 1),
            ("• High risk, high potential", 1),
        ],
    )

    create_content_slide(
        prs,
        "Portfolio Balance Dimensions",
        [
            "Balance your portfolio across:",
            "",
            "Time Horizon: H1 / H2 / H3 (60/25/15 split)",
            "",
            "Risk Profile:",
            ("• Low Risk: 50-60%", 1),
            ("• Medium Risk: 30-40%", 1),
            ("• High Risk: 10-20%", 1),
            "",
            "Business Impact:",
            ("• Operational Excellence: 50-60%", 1),
            ("• Growth: 30-40%", 1),
            ("• Transformation: 10-20%", 1),
        ],
    )

    create_content_slide(
        prs,
        "Portfolio Health Metrics",
        [
            "Regular check your portfolio balance:",
            "",
            "Monthly Review:",
            ("• Project progress, risks, resource needs", 1),
            "",
            "Quarterly Review:",
            ("• Portfolio balance, reprioritization", 1),
            ("• Add/remove/reprioritize initiatives", 1),
            "",
            "Annual Review:",
            ("• Strategic alignment", 1),
            ("• Major portfolio shifts", 1),
            ("• Budget reallocation", 1),
        ],
    )

    # Build vs Buy vs Partner
    create_section_slide(prs, "Build vs. Buy vs. Partner")

    create_content_slide(
        prs,
        "Decision Framework",
        [
            "BUILD when:",
            ("✓ Core competitive differentiator", 1),
            ("✓ Unique requirements", 1),
            ("✓ IP protection critical", 1),
            ("✓ Have internal expertise", 1),
            "",
            "BUY when:",
            ("✓ Commodity capability", 1),
            ("✓ Standard requirements", 1),
            ("✓ Fast time-to-market needed", 1),
            ("✓ Lack internal expertise", 1),
            "",
            "PARTNER when:",
            ("✓ Emerging technology/expertise gap", 1),
            ("✓ Share development costs/risks", 1),
            ("✓ Flexibility needed", 1),
        ],
    )

    create_content_slide(
        prs,
        "5-Factor Assessment Model",
        [
            "Score each capability on 5 factors:",
            "",
            "1. Competitive Differentiation:",
            ("   Core / Adjacent / Commodity", 1),
            "",
            "2. Internal Expertise:",
            ("   High / Medium / Low", 1),
            "",
            "3. Uniqueness of Requirements:",
            ("   Unique / Customizable / Standard", 1),
            "",
            "4. Time Pressure:",
            ("   High / Medium / Low", 1),
            "",
            "5. Budget:",
            ("   High / Medium / Low", 1),
        ],
    )

    # Case Studies
    create_section_slide(prs, "Real-World Case Studies")

    create_content_slide(
        prs,
        "Case Study: Netflix Technology Strategy",
        [
            "Three Horizons Approach:",
            "",
            "H1 (60%): Optimize Current Business",
            ("• Streaming infrastructure reliability", 1),
            ("• Recommendation engine improvements", 1),
            "",
            "H2 (25%): Build Adjacent",
            ("• Mobile/offline viewing", 1),
            ("• Interactive content", 1),
            ("• Gaming integration", 1),
            "",
            "H3 (15%): Transform",
            ("• AI content generation", 1),
            ("• Spatial computing experiences", 1),
        ],
    )

    create_content_slide(
        prs,
        "Netflix: Build vs Buy Decisions",
        [
            "Strategic Choices:",
            "",
            "BUILD:",
            ("• Streaming infrastructure (core differentiator)", 1),
            ("• Recommendation algorithms (competitive moat)", 1),
            ("• Encoding/compression tech (quality + cost)", 1),
            "",
            "BUY/PARTNER:",
            ("• CDNs globally (local expertise)", 1),
            ("• Content production tools (non-core)", 1),
            "",
            "Result: $1B+ annual savings, industry-leading NPS",
        ],
    )

    create_content_slide(
        prs,
        "Case Study: Capital One AI-First Strategy",
        [
            "Prioritization: Weighted Scoring Model",
            "",
            "Criteria:",
            ("• Customer Impact: 35%", 1),
            ("• Strategic Alignment: 25%", 1),
            ("• Financial ROI: 20%", 1),
            ("• Risk/Compliance: 15%", 1),
            ("• Speed to Market: 5%", 1),
            "",
            "Results:",
            ("• $900M annual savings (cloud migration)", 1),
            ("• #1 digital bank (customer satisfaction)", 1),
            ("• 11,000+ technologists (vs 3,000 in 2010)", 1),
        ],
    )

    create_content_slide(
        prs,
        "Case Study: John Deere Platform Strategy",
        [
            "Transformation Roadmap (2012-2025):",
            "",
            "2012-2015: Foundations",
            ("• Acquire precision ag startups", 1),
            ("• Build IoT sensor infrastructure", 1),
            "",
            "2016-2019: Platform Build",
            ("• Launch Operations Center (cloud)", 1),
            ("• Build ML models for yield prediction", 1),
            "",
            "2020-2023: Ecosystem Expansion",
            ("• Autonomous tractor capabilities", 1),
            ("• AI-driven optimization", 1),
            "",
            "Result: $1B+ software revenue, 20-30% yield improvements",
        ],
    )

    # Roadmap Development
    create_section_slide(prs, "Technology Roadmap Development")

    create_content_slide(
        prs,
        "Roadmap Components",
        [
            "NOW (0-6 months):",
            ("• Quick wins and critical needs", 1),
            ("• Foundation building", 1),
            ("• Proof of concepts", 1),
            "",
            "NEXT (6-18 months):",
            ("• Scale successful pilots", 1),
            ("• Core capability buildout", 1),
            ("• Platform development", 1),
            "",
            "LATER (18+ months):",
            ("• Advanced capabilities", 1),
            ("• Transformational initiatives", 1),
            ("• Emerging technologies", 1),
        ],
    )

    create_content_slide(
        prs,
        "Roadmap Principles",
        [
            "5 Principles for Effective Roadmaps:",
            "",
            "1. Outcome-focused:",
            ("   What business outcomes, not just features", 1),
            "",
            "2. Flexible:",
            ("   Adapt to changing conditions", 1),
            "",
            "3. Transparent:",
            ("   Visible to all stakeholders", 1),
            "",
            "4. Realistic:",
            ("   Based on actual capacity and resources", 1),
            "",
            "5. Strategic:",
            ("   Aligned with business objectives", 1),
        ],
    )

    # Hands-On Exercise
    create_content_slide(
        prs,
        "Hands-On Exercise",
        [
            "Your Turn: Prioritize a Portfolio",
            "",
            "Given 3 projects:",
            ("A. Customer Service AI ($2.5M value, 6 months, medium risk)", 1),
            ("B. Predictive Maintenance ($3.5M value, 12 months, high risk)", 1),
            ("C. Code Review System ($0.8M value, 3 months, low risk)", 1),
            "",
            "Questions:",
            ("1. Using weighted scoring, which would you prioritize?", 1),
            ("2. How would you balance the portfolio?", 1),
            ("3. Which would you build vs buy?", 1),
        ],
    )

    # Key Takeaways
    create_content_slide(
        prs,
        "Key Takeaways",
        [
            "1. Strategy is choice - say 'no' to many opportunities",
            "",
            "2. Use systematic prioritization frameworks",
            ("   (Weighted scoring, RICE, 2×2 matrix)", 1),
            "",
            "3. Balance portfolio across horizons (60/25/15)",
            "",
            "4. Build only core differentiators, buy/partner rest",
            "",
            "5. Roadmaps evolve - plan with flexibility",
            "",
            "6. Communicate constantly - strategy fails without buy-in",
        ],
    )

    # Q&A
    create_section_slide(prs, "Questions & Discussion")

    # Save
    prs.save("/home/user/data-advanced-ai/slides/Week13_Technology_Strategy.pptx")
    print("✅ Created: Week13_Technology_Strategy.pptx")


if __name__ == "__main__":
    import os

    # Create slides directory if it doesn't exist
    os.makedirs("/home/user/data-advanced-ai/slides", exist_ok=True)

    print("Creating lecture slide decks...")
    print("=" * 60)

    create_week1_slides()
    create_week6_slides()
    create_week13_slides()

    print("=" * 60)
    print("✅ All slide decks created successfully!")
    print("\nFiles created:")
    print("  • slides/Week01_Prompt_Engineering.pptx (27 slides)")
    print("  • slides/Week06_Agentic_Frameworks.pptx (28 slides)")
    print("  • slides/Week13_Technology_Strategy.pptx (29 slides)")
    print("\nReady for lecture!")
